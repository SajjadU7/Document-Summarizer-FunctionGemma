import sys
import re
import os
import json  # <--- Added json library
import torch
from transformers import AutoProcessor, AutoModelForCausalLM
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# --- 1. File Reading Tools (Now with Cleanup) ---
def read_file_content(file_path):
    """Reads text content and cleans up extra whitespace."""
    if not os.path.exists(file_path):
        return None
    
    ext = file_path.split('.')[-1].lower()
    text = ""

    try:
        if ext == 'pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif ext in ['docx', 'doc']:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == 'pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ['txt', 'md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            return None
    except Exception:
        return None
    
    # CLEANUP: Replaces multiple spaces/newlines with a single space
    # This fixes the "P r i n c i p l e s   o f" looking text
    clean_text = re.sub(r'\s+', ' ', text).strip()
    return clean_text[:5000] 

# --- 2. The Actual Tool Logic (Returns JSON) ---
def perform_summarization(file_path, tone="Normal"):
    content = read_file_content(file_path)
    
    # Handle errors gracefully in JSON
    if content is None:
        error_response = {
            "status": "error",
            "message": "File not found or unsupported format",
            "file_path": file_path
        }
        return json.dumps(error_response, indent=4)

    # Simulated Summary Logic
    word_count = len(content.split())
    
    # Heuristic: Define the "Intro" based on tone
    summary_intro = "Here is a summary of the document:"
    if tone.lower() == "casual":
        summary_intro = "Hey! So here's the gist of what's in the file:"
    elif tone.lower() == "formal":
        summary_intro = "The provided document comprises the following core information:"
    elif tone.lower() == "concise":
        summary_intro = "Key points:"

    # Construct the JSON Object
    response_data = {
        "status": "success",
        "file_path": file_path,
        "tone": tone,
        "meta": {
            "word_count": word_count,
            "char_count": len(content)
        },
        "summary": {
            "intro": summary_intro,
            # In a real app, 'content' below would be the output from a Generator LLM
            "preview_text": content[:500] + "..." 
        }
    }

    # Return as a JSON String
    return json.dumps(response_data, indent=4)

# --- 3. FunctionGemma Setup ---
def setup_model():
    print("Loading FunctionGemma...")
    model_id = "google/functiongemma-270m-it"
    processor = AutoProcessor.from_pretrained(model_id, device_map="auto")
    model = AutoModelForCausalLM.from_pretrained(model_id, device_map="auto", torch_dtype="auto")
    return processor, model

# --- 4. Tool Schema ---
tools_schema = [
    {
        "type": "function",
        "function": {
            "name": "summarize_document",
            "description": "Summarizes a document file with a specific tone.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string"},
                    "tone": {"type": "string", "enum": ["Casual", "Formal", "Normal", "Concise"]},
                },
                "required": ["file_path"],
            },
        },
    }
]

# --- 5. Output Parser ---
def parse_and_execute(output_text):
    match = re.search(r"<start_function_call>call:(.*?)\{(.*?)\}<end_function_call>", output_text)
    
    if match:
        func_name = match.group(1)
        args_str = match.group(2)
        
        args = {}
        string_params = re.findall(r'(\w+):<escape>(.*?)<escape>', args_str)
        for k, v in string_params:
            args[k] = v
        
        simple_params = re.findall(r'(\w+):([^<,]+)', args_str)
        for k, v in simple_params:
            if k not in args: args[k] = v

        if func_name == "summarize_document":
            return perform_summarization(args.get("file_path"), args.get("tone", "Normal"))
    
    # Fallback JSON if no function call was found
    return json.dumps({"status": "error", "message": "No function call detected."}, indent=4)

# --- Main Execution ---
if __name__ == "__main__":
    processor, model = setup_model()
    
    if len(sys.argv) > 1:
        user_prompt = sys.argv[1]
    else:
        user_prompt = input("\nENTER COMMAND: ")

    messages = [
        {"role": "developer", "content": "You are a model that can do function calling with the following functions"},
        {"role": "user", "content": user_prompt}
    ]

    inputs = processor.apply_chat_template(messages, tools=tools_schema, add_generation_prompt=True, return_dict=True, return_tensors="pt")
    inputs = {k: v.to(model.device) for k, v in inputs.items()}

    out = model.generate(**inputs, max_new_tokens=128)
    raw_output = processor.decode(out[0][len(inputs["input_ids"][0]):], skip_special_tokens=False)
    
    # This will now print pure JSON
    print(parse_and_execute(raw_output))